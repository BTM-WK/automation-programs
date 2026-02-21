#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
SPD Phase 1-A: Auto-Fetch Engine
==================================
RFP Radar v8.3 결과에서 S/A등급 공고의 첨부파일을 자동 다운로드하고
HWP→PDF 변환 → 텍스트 추출까지 수행합니다.

Usage:
  python spd_auto_fetch.py                        # 최술 리포트, S/A 등급
  python spd_auto_fetch.py --list                 # 대상 목록마 확인
  python spd_auto_fetch.py --grade S              # S등급만
  python spd_auto_fetch.py --bid 20260214001      # 특정 공고
  python spd_auto_fetch.py --report report_20260214.json  # 특정 리포트

Author: WKMG Automation (SPD System)
Version: 1.0.0
"""

import os
import sys
import json
import glob
import time
import hashlib
import argparse
import logging
import subprocess
import tempfile
from pathlib import Path
from datetime import datetime
from typing import Optional, Dict, List, Tuple

import requests

try:
    import pdfplumber
except ImportError:
    pdfplumber = None
    print("⚠️ pdfplumber 미설치: pip install pdfplumber")

# ═══════════════════════════════════════════════════════════════
# 설정
# ═══════════════════════════════════════════════════════════════

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
CONFIG_PATH = os.path.join(SCRIPT_DIR, "config.json")

DEFAULT_CONFIG = {
    "rfp_radar_data_dir": os.path.join(SCRIPT_DIR, "..", "rfp_radar", "data", "daily_reports"),
    "download_dir": os.path.join(SCRIPT_DIR, "data", "downloads"),
    "extracted_dir": os.path.join(SCRIPT_DIR, "data", "extracted"),
    "output_dir": os.path.join(SCRIPT_DIR, "data", "fetch_results"),
    "libreoffice_path": "",
    "max_file_seq": 20,
    "download_timeout": 30,
    "min_grade": "A",
    "target_sources": ["g2b_api"],
}

def load_config():
    cfg = DEFAULT_CONFIG.copy()
    if os.path.exists(CONFIG_PATH):
        with open(CONFIG_PATH, "r", encoding="utf-8") as f:
            user_cfg = json.load(f)
        if "spd" in user_cfg:
            cfg.update(user_cfg["spd"])
    # 환경변수 우선
    if os.environ.get("RFP_SERVICE_KEY"):
        cfg["service_key"] = os.environ["RFP_SERVICE_KEY"]
    return cfg

# 로깅
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%H:%M:%S"
)
log = logging.getLogger("AutoFetch")

# 등금 우선순위
GRADE_ORDER = {"S": 0, "A": 1, "B": 2, "C": 3}

# ═══════════════════════════════════════════════════════════════
# G2B 첨부파일 다운로드
# ═══════════════════════════════════════════════════════════════

# G2B 첨부파일 API — 차세대 나라장터 (2025.1.6~) + 구 API fallback
# ★ 차세대 API: /ad/BidPublicInfoService/ (R26BK... 공고번호 지원)
# ★ 구 API: /BidPublicInfoService04/ (구 공고번호만 지원 — fallback용)
G2B_FILE_URLS = {
    # ── 차세대 나라장터 공고정보 조회 API (ntceSpecDocUrl1~10에 첨부파일 포함) ──
    "servc":    "https://apis.data.go.kr/1230000/ad/BidPublicInfoService/getBidPblancListInfoServc",
    "thng":     "https://apis.data.go.kr/1230000/ad/BidPublicInfoService/getBidPblancListInfoThng",
    "cnstwk":   "https://apis.data.go.kr/1230000/ad/BidPublicInfoService/getBidPblancListInfoCnstwk",
    "frgcpt":   "https://apis.data.go.kr/1230000/ad/BidPublicInfoService/getBidPblancListInfoFrgcpt",
    "etc":      "https://apis.data.go.kr/1230000/ad/BidPublicInfoService/getBidPblancListInfoEtc",
    # ── 차세대 첨부파일 전용 API (e발주 첨부) ──
    "eorder":   "https://apis.data.go.kr/1230000/ad/BidPublicInfoService/getBidPblancListInfoEorderAtchFileInfo",
    "innov_rfp":"https://apis.data.go.kr/1230000/ad/BidPublicInfoService/getBidPblancListPPIFnlRfpIssAtchFileInfo",
}

def find_latest_report(data_dir: str) -> Optional[str]:
    """rfp_radar/data/daily_reports에서 최신 추천 리포트 찾기
    우선순위: recommend > rfp_report > report > candidates > 기타
    """
    # 1순위: recommend 파일 (S/A등급 추천 공고 — 이미 등급 필터링 완료)
    priority_patterns = [
        os.path.join(data_dir, "v83_recommend_*.json"),
        os.path.join(data_dir, "v*_recommend_*.json"),
        os.path.join(data_dir, "rfp_report_*.json"),
        os.path.join(data_dir, "report_*.json"),
    ]
    for p in priority_patterns:
        files = glob.glob(p)
        if files:
            chosen = max(files, key=os.path.getmtime)
            log.info(f"   📄 선택된 리포트: {os.path.basename(chosen)}")
            return chosen
    
    # 2순위: candidates 파일 (전체 후보 — 등급 필터 필요)
    candidates = glob.glob(os.path.join(data_dir, "v83_candidates_*.json"))
    if candidates:
        chosen = max(candidates, key=os.path.getmtime)
        log.info(f"   📄 선택된 리포트 (candidates): {os.path.basename(chosen)}")
        return chosen
    
    # 3순위: 아무 JSON
    all_files = glob.glob(os.path.join(data_dir, "*.json"))
    if not all_files:
        return None
    chosen = max(all_files, key=os.path.getmtime)
    log.info(f"   📄 선택된 리포트 (fallback): {os.path.basename(chosen)}")
    return chosen

def load_report(report_path: str) -> List[Dict]:
    """RFP Radar 리포트에서 공고 목록 로드"""
    with open(report_path, "r", encoding="utf-8") as f:
        data = json.load(f)
    
    # 다양한 리포트 형시 지원
    if isinstance(data, list):
        return data
    if isinstance(data, dict):
        for key in ["results", "bids", "items", "data"]:
            if key in data and isinstance(data[key], list):
                return data[key]
        # top-level dict with bid data
        if "bid_no" in data or "bidNtceNo" in data:
            return [data]
    return []

def filter_bids(bids: List[Dict], min_grade: str = "A", 
                target_sources: List[str] = None, 
                specific_bid: str = None) -> List[Dict]:
    """등급/소스/공고번호 기준 필터링"""
    min_grade_idx = GRADE_ORDER.get(min_grade, 1)
    filtered = []
    
    for bid in bids:
        # 특정 공고번호 지정
        if specific_bid:
            bid_no = str(bid.get("bid_no", bid.get("bidNtceNo", "")))
            if specific_bid not in bid_no:
                continue
            filtered.append(bid)
            continue
        
        # 등급 필터 (다양한 리포트 형식 지원)
        grade = bid.get("grade", bid.get("rfp_grade", ""))
        if not grade and isinstance(bid.get("score"), dict):
            grade = bid["score"].get("grade", "")
        if not grade:
            grade = "C"
        if GRADE_ORDER.get(grade, 3) > min_grade_idx:
            continue
        
        # 소스 필터
        if target_sources:
            source = bid.get("source", bid.get("data_source", ""))
            if source and not any(t in source.lower() for t in target_sources):
                continue
        
        filtered.append(bid)
    
    return filtered

def download_g2b_files(bid_no: str, config: Dict) -> List[Dict]:
    """G2B API로 공고 첨부파일 다운로드 (용역 우선, 실패 시 물품/공사/외자 순차 시도)"""
    service_key = config.get("service_key", os.environ.get("RFP_SERVICE_KEY", ""))
    if not service_key:
        log.warning(f"  ⚠️ G2B API 키 없음 — 첨부파일 다운로드 스킵")
        return []
    
    download_dir = config.get("download_dir", "data/downloads")
    bid_dir = os.path.join(download_dir, str(bid_no))
    os.makedirs(bid_dir, exist_ok=True)
    
    timeout = config.get("download_timeout", 30)
    
    # 1단계: 용역 공고정보에서 ntceSpecDocUrl 추출 (가장 확실)
    api_try_order = ["servc", "thng", "cnstwk", "frgcpt", "etc"]
    
    for api_type in api_try_order:
        api_url = G2B_FILE_URLS.get(api_type)
        downloaded = _try_download_from_bid_info(bid_no, api_url, api_type, bid_dir, service_key, timeout)
        
        if downloaded:
            log.info(f"    ✅ {api_type} 공고정보에서 {len(downloaded)}개 파일 다운로드 성공")
            return downloaded
        else:
            log.debug(f"    - {api_type}: 해당 없음")
    
    # 2단계: e발주 전용 첨부파일 API
    for api_type in ["eorder", "innov_rfp"]:
        api_url = G2B_FILE_URLS.get(api_type)
        downloaded = _try_download_from_atch_api(bid_no, api_url, api_type, bid_dir, service_key, timeout)
        
        if downloaded:
            log.info(f"    ✅ {api_type} 첨부API에서 {len(downloaded)}개 파일 다운로드 성공")
            return downloaded
        else:
            log.debug(f"    - {api_type} 첨부API: 파일 없음")
    
    # 모든 API에서 실패 — 차세대 나라장터 직접 크롤링 시도
    log.info(f"  🔄 G2B API 첨부파일 없음 → 나라장터 웹 직접 시도...")
    downloaded = _try_download_from_g2b_web(bid_no, bid_dir, timeout)
    if downloaded:
        log.info(f"    ✅ 웹 크롤링으로 {len(downloaded)}개 파일 다운로드 성공")
        return downloaded
    
    log.warning(f"  ❌ 첨부파일 다운로드 실패 (차세대+구 API + 웹 모두 실패)")
    return []


def _try_download_from_bid_info(bid_no: str, api_url: str, api_type: str,
                                bid_dir: str, service_key: str, timeout: int) -> List[Dict]:
    """공고정보 조회 API에서 ntceSpecDocUrl1~10 추출하여 파일 다운로드"""
    downloaded = []
    params = {
        "serviceKey": service_key,
        "pageNo": "1",
        "numOfRows": "10",
        "type": "json",
        "inqryDiv": "2",
        "bidNtceNo": bid_no,
    }
    try:
        resp = requests.get(api_url, params=params, timeout=timeout)
        if resp.status_code != 200:
            return []
        
        data = resp.json()
        body = data.get("response", {}).get("body", {})
        total_count = int(body.get("totalCount", 0))
        if total_count == 0:
            return []
        
        # items에서 첫 번째 아이템
        item_list = body.get("items", [])
        if isinstance(item_list, dict):
            item_list = item_list.get("item", [])
        if isinstance(item_list, dict):
            item_list = [item_list]
        if not item_list:
            return []
        
        item = item_list[0] if isinstance(item_list, list) else item_list
        
        # ntceSpecDocUrl1~10, ntceSpecFileNm1~10 추출
        for seq in range(1, 11):
            file_url = item.get(f"ntceSpecDocUrl{seq}", "")
            file_name = item.get(f"ntceSpecFileNm{seq}", "")
            
            if not file_url:
                continue
            if not file_name:
                file_name = f"ntceSpec_{seq}"
            
            file_path = os.path.join(bid_dir, file_name)
            try:
                file_resp = requests.get(file_url, timeout=timeout)
                if file_resp.status_code == 200 and len(file_resp.content) > 100:
                    with open(file_path, "wb") as f:
                        f.write(file_resp.content)
                    downloaded.append({
                        "file_name": file_name,
                        "file_path": file_path,
                        "file_size": len(file_resp.content),
                        "file_seq": seq,
                        "api_type": api_type,
                    })
                    log.info(f"    📥 [{api_type}] {file_name} ({len(file_resp.content):,}B)")
            except Exception as e:
                log.warning(f"    ⚠️ 다운로드 실패: {file_name} — {e}")
    
    except Exception as e:
        log.debug(f"    - [{api_type}] 요청 실패: {e}")
    
    return downloaded


def _try_download_from_atch_api(bid_no: str, api_url: str, api_type: str,
                                bid_dir: str, service_key: str, timeout: int) -> List[Dict]:
    """e발주/혁신장터 첨부파일 전용 API에서 다운로드"""
    downloaded = []
    params = {
        "serviceKey": service_key,
        "pageNo": "1",
        "numOfRows": "100",
        "type": "json",
        "inqryDiv": "2",
        "bidNtceNo": bid_no,
    }
    try:
        resp = requests.get(api_url, params=params, timeout=timeout)
        if resp.status_code != 200:
            return []
        
        data = resp.json()
        body = data.get("response", {}).get("body", {})
        total_count = int(body.get("totalCount", 0))
        if total_count == 0:
            return []
        
        item_list = body.get("items", [])
        if isinstance(item_list, dict):
            item_list = item_list.get("item", [])
        if isinstance(item_list, dict):
            item_list = [item_list]
        items = item_list if isinstance(item_list, list) else []
        
        log.info(f"    📋 [{api_type}] {len(items)}개 첨부파일 발견")
        
        for item in items:
            file_url = (item.get("eorderAtchFileUrl") or 
                       item.get("fileUrl") or 
                       item.get("pblancAtchFileUrl") or "")
            file_name = (item.get("eorderAtchFileNm") or 
                        item.get("fileNm") or 
                        item.get("pblancAtchFileNm") or 
                        f"file_{item.get('atchSno', '0')}")
            
            if not file_url:
                continue
            
            file_path = os.path.join(bid_dir, file_name)
            try:
                file_resp = requests.get(file_url, timeout=timeout)
                if file_resp.status_code == 200 and len(file_resp.content) > 100:
                    with open(file_path, "wb") as f:
                        f.write(file_resp.content)
                    downloaded.append({
                        "file_name": file_name,
                        "file_path": file_path,
                        "file_size": len(file_resp.content),
                        "file_seq": item.get("atchSno", 0),
                        "api_type": api_type,
                        "doc_type": item.get("eorderDocDivNm", ""),
                    })
                    log.info(f"    📥 [{api_type}] {file_name} ({len(file_resp.content):,}B)")
            except Exception as e:
                log.warning(f"    ⚠️ 다운로드 실패: {file_name} — {e}")
                
    except Exception as e:
        log.debug(f"    - [{api_type}] 요청 실패: {e}")
    
    return downloaded


def _try_download_from_g2b_web(bid_no: str, bid_dir: str, timeout: int) -> List[Dict]:
    """차세대 나라장터 웹에서 직접 첨부파일 다운로드 시도 (API 실패 시 fallback)"""
    downloaded = []
    
    try:
        # 차세대 나라장터 공고 상세 페이지에서 첨부파일 링크 추출
        detail_url = f"https://www.g2b.go.kr/link/PNPE027_01/single/?bidPbancNo={bid_no}&bidPbancOrd=00"
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36",
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
        }
        resp = requests.get(detail_url, headers=headers, timeout=timeout, allow_redirects=True)
        
        if resp.status_code != 200:
            return []
        
        html = resp.text
        
        # 첨부파일 다운로드 URL 패턴 추출 (차세대 나라장터)
        import re
        # 패턴 1: fileDownload 링크
        file_patterns = re.findall(
            r'href=["\']([^"\']*(?:fileDownload|atchFileDown|download)[^"\']*)["\']',
            html, re.IGNORECASE
        )
        # 패턴 2: 첨부파일 직접 URL
        file_patterns += re.findall(
            r'href=["\']([^"\']*\.(?:hwp|hwpx|pdf|doc|docx|xlsx|xls|zip)[^"\']*)["\']',
            html, re.IGNORECASE
        )
        
        if not file_patterns:
            return []
        
        for i, file_url in enumerate(file_patterns[:10]):  # 최대 10개
            if not file_url.startswith("http"):
                file_url = "https://www.g2b.go.kr" + file_url
            
            try:
                file_resp = requests.get(file_url, headers=headers, timeout=timeout, allow_redirects=True)
                if file_resp.status_code == 200 and len(file_resp.content) > 100:
                    # Content-Disposition에서 파일명 추출
                    cd = file_resp.headers.get("Content-Disposition", "")
                    file_name = f"web_file_{i+1}"
                    if "filename" in cd:
                        name_match = re.search(r'filename[*]?=["\']?(?:UTF-8\'\')?([^"\';\n]+)', cd)
                        if name_match:
                            from urllib.parse import unquote
                            file_name = unquote(name_match.group(1).strip())
                    
                    file_path = os.path.join(bid_dir, file_name)
                    with open(file_path, "wb") as f:
                        f.write(file_resp.content)
                    downloaded.append({
                        "file_name": file_name,
                        "file_path": file_path,
                        "file_size": len(file_resp.content),
                        "file_seq": i + 1,
                        "api_type": "web_crawl",
                    })
                    log.info(f"    📥 [웹] {file_name} ({len(file_resp.content):,}B)")
            except Exception:
                continue
                
    except Exception as e:
        log.debug(f"    웹 크롤링 실패: {e}")
    
    return downloaded

# ═══════════════════════════════════════════════════════════════
# 텍스트 추출
# ═══════════════════════════════════════════════════════════════

def extract_text_from_pdf(file_path: str) -> str:
    """PDF에서 텍스트 추출"""
    if not pdfplumber:
        return ""
    try:
        texts = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages[:50]:  # 최대 50페이지
                text = page.extract_text()
                if text:
                    texts.append(text)
        return "\n".join(texts)
    except Exception as e:
        log.warning(f"    ⚠️ PDF 추출 실패: {e}")
        return ""

def convert_hwp_to_pdf(hwp_path: str, config: Dict) -> Optional[str]:
    """HWP → PDF 변환 (LibreOffice 사용)"""
    lo_path = config.get("libreoffice_path", "")
    
    # LibreOffice 경로 자동 탐색
    if not lo_path:
        candidates = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            "/usr/bin/libreoffice",
            "/usr/bin/soffice",
        ]
        for c in candidates:
            if os.path.exists(c):
                lo_path = c
                break
    
    if not lo_path:
        log.warning("    ⚠️ LibreOffice 미설치 — HWP 변환 불가")
        return None
    
    output_dir = os.path.dirname(hwp_path)
    try:
        cmd = [lo_path, "--headless", "--convert-to", "pdf", "--outdir", output_dir, hwp_path]
        result = subprocess.run(cmd, capture_output=True, timeout=60)
        
        pdf_path = os.path.splitext(hwp_path)[0] + ".pdf"
        if os.path.exists(pdf_path):
            log.info(f"    📄 HWP→PDF 변환 완료")
            return pdf_path
    except Exception as e:
        log.warning(f"    ⚠️ HWP 변환 실패: {e}")
    
    return None

def extract_text_from_file(file_path: str, config: Dict) -> str:
    """파일 형식에 따를 텍스트 추출"""
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    
    elif ext in (".hwp", ".hwpx"):
        pdf_path = convert_hwp_to_pdf(file_path, config)
        if pdf_path:
            return extract_text_from_pdf(pdf_path)
        return ""
    
    elif ext == ".docx":
        try:
            import docx
            doc = docx.Document(file_path)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception:
            return ""
    
    elif ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
            return "\n".join(texts)
        except Exception:
            return ""
    
    elif ext == ".txt":
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        except:
            try:
                with open(file_path, "r", encoding="euc-kr") as f:
                    return f.read()
            except:
                return ""
    
    return ""

# ═══════════════════════════════════════════════════════════════
# 메인 파이프라인
# ═══════════════════════════════════════════════════════════════

def process_bid(bid: Dict, config: Dict, dry_run: bool = False) -> Dict:
    """단일 공고 처리: 다운로드 → 추출 → 결과 저장"""
    bid_no = str(bid.get("bid_no", bid.get("bidNtceNo", "unknown")))
    title = bid.get("title", bid.get("bidNtceNm", ""))
    grade = bid.get("grade", bid.get("rfp_grade", ""))
    if not grade and isinstance(bid.get("score"), dict):
        grade = bid["score"].get("grade", "?")
    agency = bid.get("agency", bid.get("dminsttNm", ""))
    budget = bid.get("budget_str", bid.get("budget", bid.get("presmptPrce", "")))
    
    log.info(f"\n{'='*60}")
    log.info(f"📋 [{grade}] {title}")
    log.info(f"   공고: {bid_no} | 기관: {agency} | 예산: {budget}")
    
    result = {
        "bid_no": bid_no,
        "title": title,
        "grade": grade,
        "agency": agency,
        "budget_str": str(budget),
        "processed_at": datetime.now().isoformat(),
        "files_downloaded": [],
        "texts_extracted": [],
        "total_text_length": 0,
        "status": "pending",
    }
    
    if dry_run:
        result["status"] = "dry_run"
        log.info(f"   🔍 DRY-RUN — 다운로드/추출 스킵")
        return result
    
    # Step 1: 첨부파일 다운로드
    log.info(f"  📥 첨부파일 다운로드 시작...")
    downloaded = download_g2b_files(bid_no, config)
    result["files_downloaded"] = downloaded
    
    if not downloaded:
        log.info(f"  ⚠️ 첨부파일 없음 — RFP 텍스트만 사용")
        result["status"] = "no_files"
        return result
    
    # Step 2: 텍스트 추출
    log.info(f"  📄 텍스트 추출 중... ({len(downloaded)}개 파일)")
    extracted_dir = config.get("extracted_dir", "data/extracted")
    os.makedirs(extracted_dir, exist_ok=True)
    
    all_texts = []
    for dl in downloaded:
        text = extract_text_from_file(dl["file_path"], config)
        if text:
            # 추출 텍스트 저장
            text_file = os.path.join(extracted_dir, f"{bid_no}_{dl['file_seq']}.txt")
            with open(text_file, "w", encoding="utf-8") as f:
                f.write(text)
            
            all_texts.append({
                "file_name": dl["file_name"],
                "text_length": len(text),
                "text_file": text_file,
                "preview": text[:200],
            })
            log.info(f"    ✅ {dl['file_name']} → {len(text):,}자 추출")
    
    result["texts_extracted"] = all_texts
    result["total_text_length"] = sum(t["text_length"] for t in all_texts)
    result["status"] = "completed" if all_texts else "extraction_failed"
    
    log.info(f"  📊 추출 완료: {len(all_texts)}개 파일, 총 {result['total_text_length']:,}자")
    
    return result

def run_auto_fetch(args):
    """메인 실행"""
    config = load_config()
    
    # 디렉토리 생성
    for d in ["download_dir", "extracted_dir", "output_dir"]:
        path = config.get(d, "")
        if path:
            os.makedirs(os.path.join(SCRIPT_DIR, path) if not os.path.isabs(path) else path, exist_ok=True)
    
    # 리포트 로드
    if args.report:
        report_path = args.report
    else:
        data_dir = config.get("rfp_radar_data_dir", "")
        if not os.path.isabs(data_dir):
            data_dir = os.path.join(SCRIPT_DIR, data_dir)
        report_path = find_latest_report(data_dir)
    
    if not report_path or not os.path.exists(report_path):
        log.error(f"❌ RFP Radar 리포트를 찾을 수 없습니다.")
        log.info(f"   기대 경로: {config.get('rfp_radar_data_dir', '?')}")
        return
    
    report_name = os.path.basename(report_path)
    log.info(f"📂 리포트: {report_name}")
    bids = load_report(report_path)
    log.info(f"   전체 공고: {len(bids)}건")
    
    # recommend 파일이면 필터 없이 전부 분석 대상 (이미 RFP Radar가 추천한 공고)
    is_recommend = "recommend" in report_name.lower()
    
    if is_recommend and not args.bid:
        filtered = bids
        log.info(f"   📌 추천 리포트 → 전체 {len(filtered)}건 분석 대상 (필터 스킵)")
    else:
        min_grade = args.grade or config.get("min_grade", "A")
        target_sources = config.get("target_sources", None)
        filtered = filter_bids(bids, min_grade, target_sources, args.bid)
        log.info(f"   대상 공고: {filtered and len(filtered) or 0}건 (등급 {min_grade}+ 필터)")
    
    if not filtered:
        log.info("✅ 대상 공고 없음. 종료.")
        return
    
    # 목록 모드
    if args.list:
        log.info(f"\n{'='*60}")
        log.info(f"📋 Auto-Fetch 대상 목록 ({len(filtered)}건)")
        log.info(f"{'='*60}")
        for i, bid in enumerate(filtered, 1):
            grade = bid.get("grade", bid.get("rfp_grade", "?"))
            title = bid.get("title", bid.get("bidNtceNm", ""))[:50]
            bid_no = bid.get("bid_no", bid.get("bidNtceNo", ""))
            log.info(f"  {i:2d}. [{grade}] {title}... ({bid_no})")
        return
    
    # 처리 실행
    results = []
    for i, bid in enumerate(filtered, 1):
        log.info(f"\n[{i}/{len(filtered)}] 처리 중...")
        result = process_bid(bid, config, dry_run=args.dry_run)
        results.append(result)
        
        if i < len(filtered):
            time.sleep(1)  # API 부하 방지
    
    # 결과 저장
    output_dir = config.get("output_dir", "data/fetch_results")
    if not os.path.isabs(output_dir):
        output_dir = os.path.join(SCRIPT_DIR, output_dir)
    os.makedirs(output_dir, exist_ok=True)
    
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_file = os.path.join(output_dir, f"fetch_{timestamp}.json")
    
    output_data = {
        "generated_at": datetime.now().isoformat(),
        "report_source": report_path,
        "total_bids": len(bids),
        "filtered_bids": len(filtered),
        "processed": len(results),
        "results": results,
    }
    
    with open(output_file, "w", encoding="utf-8") as f:
        json.dump(output_data, f, ensure_ascii=False, indent=2)
    
    # 요약
    completed = sum(1 for r in results if r["status"] == "completed")
    total_text = sum(r.get("total_text_length", 0) for r in results)
    
    log.info(f"\n{'='*60}")
    log.info(f"✅ Auto-Fetch 완료")
    log.info(f"   처리: {len(results)}건 (성공: {completed}, 파일없음: {len(results)-completed})")
    log.info(f"   추출 텍스트: {total_text:,}자")
    log.info(f"   결과: {output_file}")
    log.info(f"{'='*60}")

# ═══════════════════════════════════════════════════════════════
# CLI
# ═══════════════════════════════════════════════════════════════

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="SPD Auto-Fetch Engine")
    parser.add_argument("--report", help="특정 RFP Radar 리폨트 JSON 경로")
    parser.add_argument("--grade", choices=["S", "A", "B", "C"], help="최소 등급 필터")
    parser.add_argument("--bid", help="특정 공고번호")
    parser.add_argument("--list", action="store_true", help="대상 목록만 확인")
    parser.add_argument("--dry-run", action="store_true", help="다운로드/추출 없이 구조만 확인")
    
    args = parser.parse_args()
    run_auto_fetch(args)
