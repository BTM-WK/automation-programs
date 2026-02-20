#!/usr/bin/env python3
"""
SPD Phase 0: WKMG 과거 프로젝트 DB 인덱싱 파이프라인
=================================================
제안서/결과보고서에서 텍스트를 추출하고 ChromaDB에 벡터화하여 저장합니다.

사용법:
    1단계: pip install python-pptx pdfplumber chromadb openai tqdm
    2단계: python spd_indexer.py --scan          # 파일 목록 스캔 및 메타데이터 추출
    3단계: python spd_indexer.py --extract        # 텍스트 추출
    4단계: python spd_indexer.py --index          # ChromaDB 벡터화
    전체:  python spd_indexer.py --all            # 1~3단계 모두 실행

작성: 2026-02-13
버전: 1.0
"""

import os
import re
import json
import sys
import hashlib
import logging
from pathlib import Path
from datetime import datetime
from typing import Optional, Dict, List, Tuple
from dataclasses import dataclass, asdict

# ============================================================
# 설정
# ============================================================

# 소스 폴더
SOURCE_DIRS = [
    r"C:\Users\yso\OneDrive\Documents\1_consulting proposals",
    r"C:\Users\yso\OneDrive\Documents\2_consulting reports",
]

# 출력 폴더
OUTPUT_DIR = r"C:\Users\yso\OneDrive\Documents\GitHub\automation-programs\spd"
EXTRACTED_DIR = os.path.join(OUTPUT_DIR, "extracted_texts")
DB_DIR = os.path.join(OUTPUT_DIR, "chromadb")
INVENTORY_FILE = os.path.join(OUTPUT_DIR, "project_inventory.json")
EXTRACTION_LOG = os.path.join(OUTPUT_DIR, "extraction_log.json")

# 지원 확장자
SUPPORTED_EXTENSIONS = {".pptx", ".ppt", ".pdf", ".PDF"}

# OpenAI 설정 (임베딩용)
EMBEDDING_MODEL = "text-embedding-3-small"
EMBEDDING_DIMENSIONS = 1536

# 청킹 설정
MAX_CHUNK_SIZE = 2000  # 문자 수 기준
CHUNK_OVERLAP = 200

# ============================================================
# 로깅 설정
# ============================================================

os.makedirs(OUTPUT_DIR, exist_ok=True)

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[
        logging.StreamHandler(),
        logging.FileHandler(os.path.join(OUTPUT_DIR, "spd_indexer.log"), encoding="utf-8"),
    ],
)
logger = logging.getLogger(__name__)


# ============================================================
# 데이터 클래스
# ============================================================

@dataclass
class ProjectFile:
    """프로젝트 파일 메타데이터"""
    file_path: str
    file_name: str
    extension: str
    file_size: int
    category: str          # "제안서" or "결과보고서"
    business_type: str     # "B2G" or "B2B" or "미분류"
    sub_category: str      # "브랜딩", "전략", etc.
    project_name: str      # 파일명에서 추출한 프로젝트명
    year: Optional[int]    # 추정 연도
    client: Optional[str]  # 추정 클라이언트
    domain: str            # 역량 도메인
    file_hash: str         # 중복 체크용


# ============================================================
# Step 1: 파일 스캔 및 메타데이터 추출
# ============================================================

def scan_files() -> List[ProjectFile]:
    """모든 프로젝트 파일을 스캔하고 메타데이터를 추출합니다."""
    logger.info("=" * 60)
    logger.info("Step 1: 파일 스캔 시작")
    logger.info("=" * 60)
    
    files = []
    seen_hashes = set()
    
    for source_dir in SOURCE_DIRS:
        if not os.path.exists(source_dir):
            logger.warning(f"폴더를 찾을 수 없음: {source_dir}")
            continue
        
        for root, dirs, filenames in os.walk(source_dir):
            # 휴지통 폴더 건너뛰기
            if "휴지통" in root:
                continue
            
            for filename in filenames:
                ext = os.path.splitext(filename)[1]
                if ext.lower() not in {e.lower() for e in SUPPORTED_EXTENSIONS}:
                    continue
                
                filepath = os.path.join(root, filename)
                
                try:
                    file_size = os.path.getsize(filepath)
                except OSError:
                    continue
                
                # 파일 해시로 중복 체크 (파일명+크기 기반 빠른 해시)
                quick_hash = hashlib.md5(f"{filename}_{file_size}".encode()).hexdigest()
                if quick_hash in seen_hashes:
                    logger.debug(f"중복 건너뛰기: {filename}")
                    continue
                seen_hashes.add(quick_hash)
                
                # 메타데이터 추출
                pf = extract_metadata(filepath, filename, ext, file_size, source_dir)
                files.append(pf)
    
    logger.info(f"총 {len(files)}개 파일 발견")
    logger.info(f"  - 제안서: {sum(1 for f in files if f.category == '제안서')}개")
    logger.info(f"  - 결과보고서: {sum(1 for f in files if f.category == '결과보고서')}개")
    logger.info(f"  - B2G: {sum(1 for f in files if f.business_type == 'B2G')}개")
    logger.info(f"  - B2B: {sum(1 for f in files if f.business_type == 'B2B')}개")
    
    # 저장
    with open(INVENTORY_FILE, "w", encoding="utf-8") as f:
        json.dump([asdict(pf) for pf in files], f, ensure_ascii=False, indent=2)
    logger.info(f"인벤토리 저장: {INVENTORY_FILE}")
    
    return files


def extract_metadata(filepath, filename, ext, file_size, source_dir):
    category = "제안서" if "1_consulting proposals" in filepath else ("결과보고서" if "2_consulting reports" in filepath else "미분류")
    
    rel_path = filepath.replace(source_dir, "")
    if "\\B2G\\" in rel_path or "/B2G/" in rel_path:
        business_type = "B2G"
    elif "\\B2B\\" in rel_path or "/B2B/" in rel_path:
        business_type = "B2B"
    else:
        business_type = guess_business_type(filename)
    
    sub_category = ""
    if "\\브랜딩\\" in rel_path or "/브랜딩/" in rel_path:
        sub_category = "브랜딩"
    elif "\\전략\\" in rel_path or "/전략/" in rel_path:
        sub_category = "전략"
    
    year = extract_year(filename)
    client = extract_client(filename)
    project_name = clean_project_name(filename)
    domain = classify_domain(filename, filepath)
    file_hash = hashlib.md5(f"{filename}_{file_size}".encode()).hexdigest()
    
    return ProjectFile(
        file_path=filepath, file_name=filename, extension=ext.lower(),
        file_size=file_size, category=category, business_type=business_type,
        sub_category=sub_category, project_name=project_name, year=year,
        client=client, domain=domain, file_hash=file_hash,
    )


def extract_year(filename):
    patterns = [
        r"20[0-2]\d[01]\d[0-3]\d",
        r"20[0-2]\d[01]\d",
        r"(?:^|[_\s])20[0-2]\d(?:[_\s.]|$)",
        r"20[0-2]\d",
    ]
    for pattern in patterns:
        match = re.search(pattern, filename)
        if match:
            year_str = match.group()[:4]
            try:
                year = int(year_str)
                if 2007 <= year <= 2026:
                    return year
            except ValueError:
                continue
    
    match = re.search(r"[PR](20\d{4})", filename)
    if match:
        try:
            return int(match.group(1)[:4])
        except ValueError:
            pass
    return None


def extract_client(filename):
    known_clients = {
        "LG전자": ["LG전자", "LGE", "LG스타일러", "LG OLED", "LGMO", "LGMC"],
        "LG디스플레이": ["LGD"],
        "롯데웰푸드": ["롯데웰푸드", "롯데푸드", "롯데식품", "롯데웰"],
        "광동제약": ["광동제약", "비타500"],
        "오리온": ["오리온", "닥터유"],
        "풀무원": ["풀무원"],
        "삼양사": ["삼양사"],
        "삼성물산": ["삼성물산", "FORAIR"],
        "삼성웰스토리": ["삼성웰스토리", "웰스토리"],
        "오뚜기": ["오뚜기"],
        "SK매직": ["SK매직"],
        "농업기술실용화재단": ["농업기술실용화재단", "FACT"],
        "한국사회적기업진흥원": ["사회적기업진흥원", "사회적기업"],
        "중소기업유통센터": ["중소기업유통센터"],
        "도루코": ["도루코", "DORCO"],
        "GS리테일": ["GS리테일", "GS홈쇼핑"],
        "농심": ["농심"],
        "골프존": ["골프존"],
        "현대L&C": ["현대L&C", "현대 L&C"],
        "효성TNC": ["효성TNC", "효성"],
        "정식품": ["정식품", "그린비아"],
        "동서식품": ["동서식품"],
        "카길코리아": ["카길", "엑셀비프"],
        "한국농업기술진흥원": ["KOAT"],
        "능률교육": ["능률교육"],
        "천호식품": ["천호식품"],
        "삼화페인트": ["삼화페인트"],
        "서울우유": ["서울우유"],
        "하이트진로": ["하이트진로", "하이트맥주", "진로"],
        "매일유업": ["매일유업", "카페라떼"],
        "본가이하랑": ["본가이하랑", "다름플러스"],
        "정관장": ["정관장"],
        "올리브영": ["올리브영"],
        "샘표": ["샘표"],
        "프린텍": ["프린텍"],
    }
    for client, keywords in known_clients.items():
        for kw in keywords:
            if kw in filename:
                return client
    return None


def guess_business_type(filename):
    b2g_keywords = ["사회적기업", "사회적경제", "농업기술", "유통채널 진출", "중소기업유통센터",
                     "관광공사", "한식재단", "공영홈쇼핑", "지역브랜드", "국가브랜드",
                     "KOAT", "저탄소", "에코플리", "아리수", "환경부", "시티투어"]
    b2b_keywords = ["LG", "롯데", "삼성", "오뚜기", "SK", "풀무원", "오리온", "광동",
                    "도루코", "DORCO", "GS리테일", "골프존", "정관장", "카길", "효성"]
    for kw in b2g_keywords:
        if kw in filename:
            return "B2G"
    for kw in b2b_keywords:
        if kw in filename:
            return "B2B"
    return "미분류"


def classify_domain(filename, filepath):
    text = filename + " " + filepath
    domain_keywords = {
        "브랜드전략": ["브랜드", "Brand", "BI", "네이밍", "Naming", "포지셔닝"],
        "마케팅전략": ["마케팅", "Marketing", "홍보", "런칭", "GTM", "커뮤니케이션"],
        "상품개발": ["제품개선", "상품개발", "상품경쟁력", "제품마켓테스트", "상품기획", "상품화", "신상품"],
        "유통/판로": ["유통채널", "판로", "유통전략", "판매전략", "소싱", "진출지원"],
        "소비자조사": ["소비자", "U&A", "FGI", "FGD", "조사", "고객만족도"],
        "사회적경제": ["사회적기업", "사회적경제"],
        "농식품": ["농식품", "농업", "6차산업", "양조장", "식품경쟁력", "농특산물"],
        "지역개발": ["지역브랜드", "도시브랜드", "관광", "힐링", "특화발전", "권역"],
    }
    scores = {}
    for domain, keywords in domain_keywords.items():
        score = sum(1 for kw in keywords if kw in text)
        if score > 0:
            scores[domain] = score
    if scores:
        return max(scores, key=scores.get)
    return "기타"


def clean_project_name(filename):
    name = os.path.splitext(filename)[0]
    prefixes = [r"^\[?WKMG\]?\s*_?\s*", r"^\[?WK마케팅그룹\]?\s*_?\s*", r"^\[?KOAT\]?\s*",
                r"^★+\s*", r"^☆+\s*", r"^P\d{6}_", r"^R\d{6}_", r"^\d{6}_"]
    for prefix in prefixes:
        name = re.sub(prefix, "", name)
    suffixes = ["_최종", "_최종본", "_수정", "_발표용", "_제출용", "_송부용", "_FIN",
                "_final", "_Final", "_PT", "_clean", "_send", "(최종)", "(수정)"]
    for suffix in suffixes:
        name = name.replace(suffix, "")
    return name.strip("_ ")


# ============================================================
# Step 2: 텍스트 추출
# ============================================================

def extract_texts(files):
    logger.info("=" * 60)
    logger.info("Step 2: 텍스트 추출 시작")
    logger.info("=" * 60)
    
    os.makedirs(EXTRACTED_DIR, exist_ok=True)
    results = {}
    success_count = 0
    fail_count = 0
    
    for i, pf in enumerate(files):
        text_file = os.path.join(EXTRACTED_DIR, f"{pf.file_hash}.txt")
        
        if os.path.exists(text_file):
            with open(text_file, "r", encoding="utf-8") as f:
                text = f.read()
            if text.strip():
                results[pf.file_hash] = text
                success_count += 1
                continue
        
        try:
            text = extract_text_from_file(pf.file_path, pf.extension)
            if text and text.strip():
                with open(text_file, "w", encoding="utf-8") as f:
                    f.write(text)
                results[pf.file_hash] = text
                success_count += 1
            else:
                fail_count += 1
                logger.warning(f"빈 텍스트: {pf.file_name}")
        except Exception as e:
            fail_count += 1
            logger.error(f"추출 실패: {pf.file_name} - {e}")
        
        if (i + 1) % 20 == 0:
            logger.info(f"  진행: {i+1}/{len(files)} (성공: {success_count}, 실패: {fail_count})")
    
    logger.info(f"텍스트 추출 완료: 성공 {success_count}, 실패 {fail_count}")
    
    log_data = {"timestamp": datetime.now().isoformat(), "total": len(files), "success": success_count, "failed": fail_count}
    with open(EXTRACTION_LOG, "w", encoding="utf-8") as f:
        json.dump(log_data, f, ensure_ascii=False, indent=2)
    
    return results


def extract_text_from_file(filepath, extension):
    ext = extension.lower()
    if ext == ".pptx":
        return extract_from_pptx(filepath)
    elif ext == ".ppt":
        return extract_from_ppt(filepath)
    elif ext == ".pdf":
        return extract_from_pdf(filepath)
    return None


def extract_from_pptx(filepath):
    try:
        from pptx import Presentation
    except ImportError:
        logger.error("python-pptx 미설치")
        return None
    try:
        prs = Presentation(filepath)
        texts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_texts:
                            slide_texts.append(" | ".join(row_texts))
            if slide_texts:
                texts.append(f"[슬라이드 {slide_num}]\n" + "\n".join(slide_texts))
        return "\n\n".join(texts)
    except Exception as e:
        logger.error(f"PPTX 오류 ({os.path.basename(filepath)}): {e}")
        return None


def extract_from_ppt(filepath):
    import subprocess, tempfile
    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            cmd = ["soffice", "--headless", "--convert-to", "pptx", "--outdir", tmpdir, filepath]
            subprocess.run(cmd, capture_output=True, timeout=60)
            for f in os.listdir(tmpdir):
                if f.endswith(".pptx"):
                    return extract_from_pptx(os.path.join(tmpdir, f))
    except FileNotFoundError:
        logger.warning(f"LibreOffice 미설치 - PPT 건너뛰기: {os.path.basename(filepath)}")
    except Exception as e:
        logger.warning(f"PPT 변환 오류: {os.path.basename(filepath)} - {e}")
    return None


def extract_from_pdf(filepath):
    try:
        import pdfplumber
    except ImportError:
        logger.error("pdfplumber 미설치")
        return None
    try:
        texts = []
        with pdfplumber.open(filepath) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    texts.append(f"[페이지 {page_num}]\n{text.strip()}")
        return "\n\n".join(texts)
    except Exception as e:
        logger.error(f"PDF 오류 ({os.path.basename(filepath)}): {e}")
        return None


# ============================================================
# Step 3: 벡터화 및 ChromaDB 저장
# ============================================================

def index_to_chromadb(files, texts):
    logger.info("=" * 60)
    logger.info("Step 3: ChromaDB 인덱싱 시작")
    logger.info("=" * 60)
    
    try:
        import chromadb
        from chromadb.utils.embedding_functions import OpenAIEmbeddingFunction
    except ImportError:
        logger.error("chromadb 미설치. pip install chromadb openai")
        return
    
    api_key = os.environ.get("OPENAI_API_KEY")
    if not api_key:
        config_paths = [
            r"C:\Users\yso\OneDrive\Documents\GitHub\automation-programs\rfp_radar\config.py",
            r"C:\Users\yso\OneDrive\Documents\GitHub\automation-programs\config.py",
        ]
        for config_path in config_paths:
            if os.path.exists(config_path):
                try:
                    with open(config_path, "r", encoding="utf-8") as f:
                        content = f.read()
                    match = re.search(r'OPENAI_API_KEY\s*=\s*["\']([^"\']+)["\']', content)
                    if match:
                        api_key = match.group(1)
                        break
                except Exception:
                    continue
    
    if not api_key:
        logger.error("OPENAI_API_KEY를 찾을 수 없습니다!")
        return
    
    os.makedirs(DB_DIR, exist_ok=True)
    client = chromadb.PersistentClient(path=DB_DIR)
    embedding_fn = OpenAIEmbeddingFunction(api_key=api_key, model_name=EMBEDDING_MODEL)
    
    collection_name = "wkmg_projects"
    try:
        client.delete_collection(collection_name)
    except Exception:
        pass
    
    collection = client.create_collection(
        name=collection_name, embedding_function=embedding_fn,
        metadata={"description": "WKMG 과거 프로젝트 제안서 및 결과보고서"},
    )
    
    total_chunks = 0
    batch_ids, batch_documents, batch_metadatas = [], [], []
    BATCH_SIZE = 50
    
    for pf in files:
        if pf.file_hash not in texts:
            continue
        text = texts[pf.file_hash]
        chunks = create_chunks(text, MAX_CHUNK_SIZE, CHUNK_OVERLAP)
        
        for i, chunk in enumerate(chunks):
            chunk_id = f"{pf.file_hash}_{i:04d}"
            metadata = {
                "file_name": pf.file_name, "category": pf.category,
                "business_type": pf.business_type, "sub_category": pf.sub_category,
                "project_name": pf.project_name, "year": pf.year or 0,
                "client": pf.client or "미확인", "domain": pf.domain,
                "chunk_index": i, "total_chunks": len(chunks),
            }
            batch_ids.append(chunk_id)
            batch_documents.append(chunk)
            batch_metadatas.append(metadata)
            total_chunks += 1
            
            if len(batch_ids) >= BATCH_SIZE:
                try:
                    collection.add(ids=batch_ids, documents=batch_documents, metadatas=batch_metadatas)
                except Exception as e:
                    logger.error(f"배치 인덱싱 오류: {e}")
                batch_ids, batch_documents, batch_metadatas = [], [], []
    
    if batch_ids:
        try:
            collection.add(ids=batch_ids, documents=batch_documents, metadatas=batch_metadatas)
        except Exception as e:
            logger.error(f"마지막 배치 오류: {e}")
    
    logger.info(f"인덱싱 완료: {total_chunks}개 청크 (원본 {len(texts)}개 파일)")
    logger.info(f"ChromaDB 저장 위치: {DB_DIR}")


def create_chunks(text, max_size=2000, overlap=200):
    if len(text) <= max_size:
        return [text]
    chunks = []
    sections = re.split(r"\n\n(?=\[(?:슬라이드|페이지) \d+\])", text)
    current_chunk = ""
    for section in sections:
        if len(current_chunk) + len(section) <= max_size:
            current_chunk += ("\n\n" if current_chunk else "") + section
        else:
            if current_chunk:
                chunks.append(current_chunk)
            if len(section) > max_size:
                sentences = re.split(r"(?<=[.!?。])\s+", section)
                current_chunk = ""
                for sentence in sentences:
                    if len(current_chunk) + len(sentence) <= max_size:
                        current_chunk += (" " if current_chunk else "") + sentence
                    else:
                        if current_chunk:
                            chunks.append(current_chunk)
                        current_chunk = sentence
            else:
                current_chunk = section
    if current_chunk:
        chunks.append(current_chunk)
    return chunks


# ============================================================
# 검색 테스트
# ============================================================

def test_search(query, n_results=5):
    try:
        import chromadb
        from chromadb.utils.embedding_functions import OpenAIEmbeddingFunction
    except ImportError:
        logger.error("chromadb 미설치")
        return
    
    api_key = os.environ.get("OPENAI_API_KEY")
    if not api_key:
        config_path = r"C:\Users\yso\OneDrive\Documents\GitHub\automation-programs\rfp_radar\config.py"
        if os.path.exists(config_path):
            with open(config_path, "r", encoding="utf-8") as f:
                content = f.read()
            match = re.search(r'OPENAI_API_KEY\s*=\s*["\']([^"\']+)["\']', content)
            if match:
                api_key = match.group(1)
    
    if not api_key:
        logger.error("OPENAI_API_KEY 필요")
        return
    
    client = chromadb.PersistentClient(path=DB_DIR)
    embedding_fn = OpenAIEmbeddingFunction(api_key=api_key, model_name=EMBEDDING_MODEL)
    collection = client.get_collection(name="wkmg_projects", embedding_function=embedding_fn)
    
    results = collection.query(query_texts=[query], n_results=n_results)
    
    print(f"\n[검색어] '{query}'")
    print("=" * 60)
    for i, (doc, metadata, distance) in enumerate(zip(
        results["documents"][0], results["metadatas"][0], results["distances"][0])):
        print(f"\n[{i+1}] 유사도: {1 - distance:.3f}")
        print(f"    프로젝트: {metadata['project_name']}")
        print(f"    구분: {metadata['category']} | {metadata['business_type']} | {metadata['domain']}")
        print(f"    연도: {metadata['year']} | 클라이언트: {metadata['client']}")
        print(f"    파일: {metadata['file_name']}")
        print(f"    내용 미리보기: {doc[:200]}...")


# ============================================================
# 메인
# ============================================================

def main():
    import argparse
    parser = argparse.ArgumentParser(description="SPD 과거 프로젝트 DB 인덱서")
    parser.add_argument("--scan", action="store_true", help="파일 스캔")
    parser.add_argument("--extract", action="store_true", help="텍스트 추출")
    parser.add_argument("--index", action="store_true", help="ChromaDB 벡터화")
    parser.add_argument("--all", action="store_true", help="전체 실행")
    parser.add_argument("--search", type=str, help="검색 테스트")
    parser.add_argument("--top", type=int, default=5, help="검색 결과 수")
    args = parser.parse_args()
    
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    
    if args.search:
        test_search(args.search, args.top)
        return
    
    if args.all or args.scan:
        files = scan_files()
    else:
        if os.path.exists(INVENTORY_FILE):
            with open(INVENTORY_FILE, "r", encoding="utf-8") as f:
                data = json.load(f)
            files = [ProjectFile(**d) for d in data]
            logger.info(f"기존 인벤토리 로드: {len(files)}개 파일")
        else:
            logger.error("인벤토리 없음. --scan 먼저 실행")
            return
    
    if args.all or args.extract:
        texts = extract_texts(files)
    else:
        texts = {}
        if os.path.exists(EXTRACTED_DIR):
            for txt_file in os.listdir(EXTRACTED_DIR):
                if txt_file.endswith(".txt"):
                    hash_id = txt_file.replace(".txt", "")
                    with open(os.path.join(EXTRACTED_DIR, txt_file), "r", encoding="utf-8") as f:
                        texts[hash_id] = f.read()
        logger.info(f"기존 추출 텍스트 로드: {len(texts)}개")
    
    if args.all or args.index:
        index_to_chromadb(files, texts)
    
    if not any([args.scan, args.extract, args.index, args.all, args.search]):
        parser.print_help()
        print('\n예시:')
        print('  python spd_indexer.py --all          # 전체 실행')
        print('  python spd_indexer.py --search "사회적기업 유통채널 진출"')


if __name__ == "__main__":
    main()
